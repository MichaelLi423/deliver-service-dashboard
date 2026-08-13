#!/usr/bin/env python3
"""Generate the 34-page Vibe Coding editorial training deck.

One 1600×900 content/layout model drives editable PPTX, fixed-render PDF and
optional PNG previews. The visual language is an annotated development dossier:
paper, terminal transcripts, document excerpts, review marks and evidence logs.
"""

from __future__ import annotations

import math
import os
from dataclasses import dataclass, field
from io import BytesIO
from pathlib import Path
from typing import Any, Iterable, cast

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt
from reportlab.lib.utils import ImageReader
from reportlab.pdfgen import canvas

W, H = 1600, 900
PPT_W, PPT_H = 13.333333, 7.5

PAPER = "F2F0EA"
INK = "111820"
SLATE = "56616D"
BLUE = "176B87"
GREEN = "4F7A57"
RED = "B84A45"
YELLOW = "D7B35C"
WHITE = "FBFAF6"
RULE = "C8C5BC"
SOFT = "E5E2DA"
TERMINAL = "17212A"
TERMINAL_TEXT = "E8EEE9"
CODE_MUTED = "9BAAA5"

MARGIN_X = 96
COL_W, GUTTER = 88, 32
CONTENT_TOP, CONTENT_BOTTOM = 248, 792

TITLES = [
    "封面", "学完能做什么", "模型 / Agent / tool / session", "受控开发", "安全红线",
    "开始前：仓库与 Git", "OpenCode 是什么", "首次会话·步骤 1：启动与状态",
    "首次会话·步骤 2：只读理解", "首次会话·步骤 3：模式与权限",
    "首次会话·步骤 4：diff / undo / 验证", "工具选择地图", "skills：标准 vs 技能包",
    "mattpocock/skills 是什么、为什么", "安装与审阅",
    "grill / tdd / diagnose / code-review 用法", "何时单 Agent",
    "oh-my-opencode-slim 是什么、为什么", "安装 / 配置 / 权限", "agents 分工与最小用法",
    "OpenSpec 是什么、为什么", "目录与 artifacts", "CLI vs 聊天命令",
    "propose → apply → verify → archive", "MCP 边界", "TDD vs diagnosis", "四层证据",
    "演示原始需求", "grill 澄清", "专用 training change / spec", "Red 与 preset-failure",
    "Green / review / verify", "日常落地 / 风险适配", "Q&A / 收尾",
]

SECTIONS = {
    range(1, 6): "概念与安全",
    range(6, 12): "第一次会话",
    range(12, 26): "工具教程",
    range(26, 28): "开发方法",
    range(28, 33): "项目演示",
    range(33, 34): "日常落地",
    range(34, 35): "问答",
}


@dataclass
class Element:
    kind: str
    x: float
    y: float
    w: float
    h: float
    props: dict[str, Any] = field(default_factory=dict)


@dataclass
class Slide:
    number: int
    title: str
    dark: bool = False
    elements: list[Element] = field(default_factory=list)

    def add(self, kind: str, x: float, y: float, w: float, h: float, **props: Any) -> None:
        self.elements.append(Element(kind, x, y, w, h, props))


def cx(col: int) -> int:
    return MARGIN_X + col * (COL_W + GUTTER)


def cw(cols: int) -> int:
    return cols * COL_W + (cols - 1) * GUTTER


def section_for(page: int) -> str:
    return next(name for pages, name in SECTIONS.items() if page in pages)


def rect(s: Slide, x: float, y: float, w: float, h: float, fill: str,
         stroke: str | None = None, sw: float = 1, radius: float = 0,
         group: str | None = None) -> None:
    s.add("rect", x, y, w, h, fill=fill, stroke=stroke, sw=sw, radius=radius, group=group)


def line(s: Slide, x1: float, y1: float, x2: float, y2: float, color: str = RULE,
         sw: float = 2, arrow: bool = False) -> None:
    s.add("line", x1, y1, x2 - x1, y2 - y1, color=color, sw=sw, arrow=arrow)


def text(s: Slide, x: float, y: float, w: float, h: float, value: str, size: float,
         color: str | None = None, weight: str = "regular", align: str = "left",
         valign: str = "top", font: str = "cjk", leading: float = 1.18,
         group: str | None = None) -> None:
    # All text starts on the 8 px baseline grid; inner document/terminal offsets
    # are snapped here rather than hand-tuned per slide.
    y = round(y / 8) * 8
    s.add("text", x, y, w, h, text=value, size=size,
          color=color or (WHITE if s.dark else INK), weight=weight, align=align,
          valign=valign, font=font, leading=leading, group=group)


def page_base(number: int, title: str | None = None, dark: bool = False) -> Slide:
    shown = title if title is not None else TITLES[number - 1]
    s = Slide(number, TITLES[number - 1], dark)
    rect(s, 0, 0, W, H, INK if dark else PAPER)
    if number not in (1, 34):
        text(s, MARGIN_X, 72, cw(2), 24, f"{number:02d}", 16,
             YELLOW if dark else BLUE, "bold", font="latin")
        text(s, cx(1), 72, cw(5), 24, section_for(number), 16,
             WHITE if dark else SLATE, "regular")
        text(s, MARGIN_X, 112, cw(10), 96, shown, 44,
             WHITE if dark else INK, "bold", leading=1.05)
    return s


def finish(s: Slide) -> Slide:
    if s.number not in (1, 34):
        line(s, MARGIN_X, 840, 1504, 840, "33414A" if s.dark else RULE, 1)
        text(s, 1456, 848, 48, 20, str(s.number), 13,
             "91A0A5" if s.dark else SLATE, align="right", font="latin")
    return s


def label(s: Slide, x: float, y: float, value: str, color: str = BLUE) -> None:
    text(s, x, y, 360, 28, value, 16, color, "bold")


def quote_mark(s: Slide, x: float, y: float, h: float, color: str = BLUE) -> None:
    rect(s, x, y, 6, h, color)


def ruled_note(s: Slide, x: float, y: float, w: float, h: float, heading: str,
               body: str, accent: str = YELLOW) -> None:
    rect(s, x, y, w, h, "F8F4E8", RULE, 1)
    rect(s, x, y, w, 10, accent)
    text(s, x + 32, y + 32, w - 64, 36, heading, 18, SLATE, "bold")
    text(s, x + 32, y + 88, w - 64, h - 112, body, 28, INK, "bold", leading=1.25)


def terminal(s: Slide, x: float, y: float, w: float, h: float, title: str,
             lines: list[tuple[str, str]], group: str | None = None) -> None:
    rect(s, x, y, w, h, TERMINAL, "2B3942", 1, 4, group)
    rect(s, x, y, w, 44, "222E37")
    text(s, x + 24, y + 11, w - 48, 24, title, 15, CODE_MUTED, "medium", font="latin")
    yy = y + 70
    for value, color in lines:
        text(s, x + 28, yy, w - 56, 34, value, 19, color, "regular", font="mono")
        yy += 38


def doc_excerpt(s: Slide, x: float, y: float, w: float, h: float, path: str,
                heading: str, body: list[str], accent: str = BLUE) -> None:
    rect(s, x + 10, y + 12, w, h, SOFT)
    rect(s, x, y, w, h, WHITE, RULE, 1)
    text(s, x + 32, y + 24, w - 64, 24, path, 15, SLATE, font="mono")
    line(s, x + 32, y + 60, x + w - 32, y + 60, RULE, 1)
    text(s, x + 32, y + 84, w - 64, 48, heading, 25, INK, "bold")
    yy = y + 152
    for item in body:
        rect(s, x + 32, yy + 8, 8, 8, accent)
        text(s, x + 56, yy, w - 88, 54, item, 20, INK, leading=1.2)
        yy += 72


def step_index(s: Slide, step: int, caption: str) -> None:
    text(s, cx(9), 72, cw(3), 28, f"STEP {step:02d} / 04", 18,
         YELLOW if s.dark else BLUE, "bold", align="right", font="latin")
    text(s, cx(9), 104, cw(3), 28, "静态会话回放", 16,
         WHITE if s.dark else SLATE, align="right")
    label(s, MARGIN_X, 216, caption, YELLOW if s.dark else BLUE)


def two_columns(s: Slide, left_head: str, left_lines: list[str], right_head: str,
                right_lines: list[str], left_color: str = BLUE,
                right_color: str = INK) -> None:
    lx, rx = cx(0), cx(6)
    line(s, cx(6) - 16, 256, cx(6) - 16, 760, RULE, 1)
    label(s, lx, 264, left_head, left_color)
    label(s, rx, 264, right_head, right_color)
    for x, items, color in ((lx, left_lines, left_color), (rx, right_lines, right_color)):
        yy = 336
        for item in items:
            text(s, x, yy, cw(5), 64, item, 24, INK, "medium", leading=1.2)
            line(s, x, yy + 72, x + cw(5), yy + 72, RULE, 1)
            yy += 112


def make_slides() -> list[Slide]:
    out: list[Slide] = []

    # 01 — editorial cover
    s = page_base(1, dark=True)
    text(s, 96, 88, cw(5), 30, "技术混合受众 · 全部静态回放", 17, YELLOW, "bold")
    text(s, 96, 184, cw(8), 108, "Vibe Coding", 62, WHITE, "bold", font="latin")
    text(s, 96, 304, cw(8), 72, "从需求到证据的受控开发闭环", 34, WHITE, "medium")
    line(s, 96, 424, cx(8) - 32, 424, "53616A", 1)
    text(s, 96, 472, cw(6), 64, "一次真实工作过程，逐页拆开。", 28, "CBD3D0", "regular")
    notes = ["90 分钟", "34 页", "不跟做", "无考核"]
    for i, item in enumerate(notes):
        x = cx(8) + (i % 2) * (cw(2) + 32)
        y = 232 + (i // 2) * 144
        text(s, x, y, cw(2), 40, item, 24, WHITE, "bold")
        line(s, x, y + 56, x + cw(2), y + 56, YELLOW if i == 0 else "53616A", 3 if i == 0 else 1)
    text(s, 96, 824, cw(8), 24, "现场不调用 Agent", 15, "91A0A5")
    out.append(s)

    # 02
    s = page_base(2)
    text(s, cx(0), 264, cw(3), 40, "01", 20, BLUE, "bold", font="latin")
    text(s, cx(0), 320, cw(3), 144, "判断\nAI 该做多少", 31, INK, "bold")
    text(s, cx(4), 264, cw(3), 40, "02", 20, BLUE, "bold", font="latin")
    text(s, cx(4), 320, cw(3), 144, "开一次\n受控会话", 31, INK, "bold")
    text(s, cx(8), 264, cw(4), 40, "03", 20, BLUE, "bold", font="latin")
    text(s, cx(8), 320, cw(4), 144, "用四层证据\n确认做对了", 31, INK, "bold")
    line(s, 96, 560, 1504, 560, RULE, 1)
    quote_mark(s, 96, 624, 56, BLUE)
    text(s, 128, 616, cw(9), 64, "现场不调用 Agent，全部内容为静态回放。", 26, INK, "medium")
    out.append(finish(s))

    # 03
    s = page_base(3)
    words = [("模型", "推理引擎"), ("Agent", "执行主体"), ("tool", "可调用能力"), ("session", "一次工作容器")]
    for i, (head, sub) in enumerate(words):
        y = 264 + i * 104
        text(s, cx(0), y, cw(2), 40, head, 28, BLUE if i == 3 else INK, "bold", font="latin" if i else "cjk")
        line(s, cx(2), y + 18, cx(4) - 32, y + 18, RULE, 1)
        text(s, cx(4), y, cw(4), 40, sub, 24, INK, "medium")
    ruled_note(s, cx(8), 264, cw(4), 416, "SESSION 边界", "工作目录\n上下文\n权限", BLUE)
    text(s, cx(8), 712, cw(4), 40, "会话是本课的操作单元", 22, BLUE, "bold")
    out.append(finish(s))

    # 04
    s = page_base(4)
    text(s, cx(0), 264, cw(8), 64, "人定义四件事", 28, BLUE, "bold")
    items = ["目标", "约束", "决策", "验收证据"]
    for i, item in enumerate(items):
        x = cx(i * 2)
        text(s, x, 352, cw(2), 64, f"{i+1:02d}", 18, SLATE, "bold", font="latin")
        text(s, x, 424, cw(2), 64, item, 28, INK, "bold")
        line(s, x, 504, x + cw(2), 504, BLUE, 3)
    quote_mark(s, cx(8), 264, 200, GREEN)
    text(s, cx(8) + 32, 264, cw(4) - 32, 128, "AI 辅助\n分析与实现", 34, INK, "bold")
    line(s, cx(8) + 32, 472, 1504, 472, RULE, 1)
    text(s, cx(8) + 32, 512, cw(4) - 32, 80, "不是\n“凭感觉接受输出”", 22, RED, "bold")
    out.append(finish(s))

    # 05
    s = page_base(5, dark=True)
    rules = [("01", "不上传敏感信息"), ("02", "不展示密钥"), ("03", "AI 无生产权限"), ("04", "高风险操作人工确认")]
    for i, (n, item) in enumerate(rules):
        y = 256 + i * 112
        text(s, cx(0), y, cw(1), 40, n, 18, YELLOW, "bold", font="latin")
        text(s, cx(1), y - 8, cw(9), 56, item, 31, WHITE, "bold")
        line(s, cx(1), y + 64, 1504, y + 64, "3C4951", 1)
    text(s, cx(8), 736, cw(4), 32, "不确定，就停下确认。", 22, "D9C982", "bold", align="right")
    out.append(finish(s))

    # 06
    s = page_base(6)
    terminal(s, cx(0), 264, cw(7), 392, "repository — clean starting point", [
        ("$ git status --short", TERMINAL_TEXT),
        ("# no output", CODE_MUTED),
        ("$ git branch --show-current", TERMINAL_TEXT),
        ("training/vibe-coding-reminder-window", "B9D7C0"),
    ])
    label(s, cx(8), 272, "开始前只确认两件事")
    text(s, cx(8), 336, cw(4), 56, "仓库干净", 29, INK, "bold")
    text(s, cx(8), 432, cw(4), 56, "分支明确", 29, INK, "bold")
    line(s, cx(8), 528, 1504, 528, RULE, 1)
    text(s, cx(8), 568, cw(4), 112, "干净、可回退的起点\n才值得信任。", 24, SLATE, "medium")
    out.append(finish(s))

    # 07
    s = page_base(7)
    label(s, cx(0), 256, "一次可检查的编码会话")
    text(s, cx(0), 312, cw(5), 112, "模型 + 工具 + 权限", 34, INK, "bold")
    text(s, cx(0), 448, cw(5), 72, "价值不在打字快，\n而在每一步都能检查。", 24, SLATE, "medium")
    rect(s, cx(6), 256, cw(6), 456, "E7E9E5", RULE, 1)
    text(s, cx(6) + 24, 280, cw(6) - 48, 28, "OpenCode · opencode.ai", 17, BLUE, "bold", font="latin")
    areas = [("对话", 336, 144), ("文件差异", 496, 96), ("终端", 608, 72)]
    for name, y, h in areas:
        rect(s, cx(6) + 24, y, cw(6) - 48, h, WHITE, RULE, 1)
        text(s, cx(6) + 48, y + 20, 240, 32, name, 20, INK, "bold")
    out.append(finish(s))

    # 08
    s = page_base(8)
    step_index(s, 1, "先看状态，再看任务")
    terminal(s, cx(0), 264, cw(8), 464, "session status — static replay", [
        ("model       selected / visible", TERMINAL_TEXT),
        ("workspace   /project/repository", TERMINAL_TEXT),
        ("tools       read · search · terminal", TERMINAL_TEXT),
        ("permissions  allow / ask by rule", "E1C980"),
    ])
    label(s, cx(9), 280, "状态回答")
    text(s, cx(9), 336, cw(3), 152, "我在哪？\n能看什么？\n能做什么？", 27, INK, "bold")
    line(s, cx(9), 544, 1504, 544, RULE, 1)
    text(s, cx(9), 584, cw(3), 96, "状态是会话的\n边界说明书。", 22, BLUE, "bold")
    out.append(finish(s))

    # 09
    s = page_base(9)
    step_index(s, 2, "先验证理解，再谈修改")
    rect(s, cx(0), 264, cw(7), 456, WHITE, RULE, 1)
    text(s, cx(0) + 32, 296, cw(7) - 64, 28, "只读请求 · 静态回放", 16, BLUE, "bold")
    text(s, cx(0) + 32, 352, cw(7) - 64, 112,
         "读取仓库，回答：\n项目是什么？如何组织？如何构建？", 26, INK, "bold")
    line(s, cx(0) + 32, 504, cx(7) - 32, 504, RULE, 1)
    text(s, cx(0) + 32, 544, cw(7) - 64, 88,
         "限制：只读，不修改任何文件。", 22, RED, "bold")
    quote_mark(s, cx(8), 288, 224, BLUE)
    text(s, cx(8) + 32, 288, cw(4) - 32, 160,
         "理解错了，\n后面一切免谈。", 32, INK, "bold")
    text(s, cx(8) + 32, 544, cw(4) - 32, 96,
         "回答先与仓库事实核对。", 21, SLATE)
    out.append(finish(s))

    # 10
    s = page_base(10)
    step_index(s, 3, "模式与权限：allow / ask / deny")
    doc_excerpt(s, cx(0), 264, cw(7), 448, "session / agent + permission", "本回放显式切到 Plan", [
        "默认主 Agent：Build",
        "默认多数 permission：allow",
        "Plan 默认编辑 / bash：ask",
    ])
    rect(s, cx(8), 304, cw(4), 240, "F7EDE8", RED, 2)
    text(s, cx(8) + 32, 336, cw(4) - 64, 28, "PERMISSION REQUEST", 16, RED, "bold", font="latin")
    text(s, cx(8) + 32, 392, cw(4) - 64, 72, "这次操作命中 ask", 26, INK, "bold")
    text(s, cx(8) + 32, 488, cw(4) - 64, 32, "看清请求，再批准", 20, RED, "bold")
    text(s, cx(8), 600, cw(4), 96, "allow 直接放行\nask 请求批准 · deny 直接拒绝", 20, SLATE, "medium")
    out.append(finish(s))

    # 11
    s = page_base(11)
    step_index(s, 4, "改完以后：把关、回退、验证")
    blocks = [
        (cx(0), "DIFF", ["范围有没有夹带？", "+ intended change", "- unrelated edit"], BLUE),
        (cx(4), "/UNDO · /REDO", ["撤销上一条用户消息", "含后续回复与改动", "可连续撤销；redo 恢复"], YELLOW),
        (cx(8), "VERIFY", ["tests", "typecheck", "build / E2E"], GREEN),
    ]
    for x, head, rows, color in blocks:
        text(s, x, 280, cw(4), 36, head, 20, color, "bold", font="latin")
        line(s, x, 328, x + cw(4), 328, color, 3)
        for j, row in enumerate(rows):
            text(s, x, 368 + j * 88, cw(4), 48, row, 22, INK, "medium")
    text(s, cx(0), 688, cw(12), 48, "diff 把关 · /undo 可连续回退 · /redo 恢复 · 跨会话 / 版本点用 Git", 22, INK, "bold")
    out.append(finish(s))

    # 12 — quiet section map
    s = page_base(12, dark=True)
    text(s, 96, 264, cw(7), 104, "工具不是一张品牌清单，\n而是四种不同职责。", 38, WHITE, "bold")
    rows = [("skills", "固化工作流"), ("编排插件", "复杂任务分角色"), ("OpenSpec", "组织规范与证据"), ("MCP", "只处理外部连接边界")]
    for i, (a, b) in enumerate(rows):
        y = 448 + i * 72
        text(s, cx(0), y, cw(3), 32, a, 22, YELLOW if i == 3 else WHITE, "bold", font="latin" if a in ("skills", "OpenSpec", "MCP") else "cjk")
        text(s, cx(4), y, cw(5), 32, b, 21, "BAC4C5")
    text(s, cx(9), 600, cw(3), 96, "MCP：本项目无配置\n本课不演示", 19, "D9C982", "bold", align="right")
    out.append(finish(s))

    # 13
    s = page_base(13)
    two_columns(s, "开放标准 / 生态", ["Agent Skills", "skills.sh", "定义如何写与分发"],
                "具体技能包", ["mattpocock/skills", "个人维护", "不是官方标准"], BLUE, RED)
    text(s, cx(0), 720, cw(12), 40, "标准定义形式；技能包提供具体内容。", 24, INK, "bold")
    out.append(finish(s))

    # 14
    s = page_base(14)
    label(s, cx(0), 256, "个人技能集合 · 可复用工作流")
    doc_excerpt(s, cx(0), 312, cw(5), 360, "skills/grill-with-docs/SKILL.md",
                "先读文档，再追问", [
                    "把模糊词变成可验证问题",
                    "把确认结果写回约定",
                ], BLUE)
    skills = [("grill-with-docs", "把歧义问成约束"), ("tdd", "先红后绿"),
              ("diagnosing-bugs", "先证据后改码"), ("code-review", "检查实现质量")]
    for i, (name, why) in enumerate(skills):
        x = cx(6 + (i % 2) * 3); y = 312 + (i // 2) * 160
        text(s, x, y, cw(3), 36, name, 19, BLUE, "bold", font="latin")
        text(s, x, y + 52, cw(3), 56, why, 22, INK, "bold")
        line(s, x, y + 120, x + cw(3), y + 120, RULE, 1)
    text(s, cx(6), 672, cw(6), 40, "四个技能覆盖受控开发的四个环节。", 21, SLATE, "medium")
    out.append(finish(s))

    # 15
    s = page_base(15)
    terminal(s, cx(0), 272, cw(8), 144, "示例命令 · 课堂不执行", [
        ("$ npx skills@latest add mattpocock/skills", TERMINAL_TEXT),
    ])
    doc_excerpt(s, cx(0), 464, cw(8), 256, "review / installed files", "启用前审阅", [
        "看清写入了哪些文件",
        "确认内容、setup 与权限边界",
    ])
    label(s, cx(9), 280, "三条纪律", RED)
    text(s, cx(9), 336, cw(3), 256,
         "安装时选择\nsetup-matt-pocock-skills\n\n每仓库运行一次\n/setup-matt-pocock-skills\n\n插件安装与 skills.sh 二选一", 17, INK, "bold", leading=1.28)
    text(s, cx(9), 656, cw(3), 48, "版本以官网 / 上游为准", 18, SLATE)
    out.append(finish(s))

    # 16
    s = page_base(16)
    rows = [("grill", "先文档，后提问"), ("tdd", "先写失败测试"),
            ("diagnose", "先收集证据，再改码"), ("code-review", "看安全、清晰、可维护")]
    for i, (name, usage) in enumerate(rows):
        y = 264 + i * 112
        text(s, cx(0), y, cw(3), 42, name, 22, BLUE, "bold", font="latin")
        text(s, cx(3), y, cw(7), 42, usage, 26, INK, "bold")
        text(s, cx(11), y, cw(1), 32, f"{i+1:02d}", 16, SLATE, "bold", align="right", font="latin")
        line(s, cx(0), y + 64, 1504, y + 64, RULE, 1)
    text(s, cx(0), 736, cw(12), 32, "四句用法，就是后半场演示的骨架。", 22, GREEN, "bold")
    out.append(finish(s))

    # 17
    s = page_base(17)
    two_columns(s, "一个 Agent 就够", ["边界清晰", "串行依赖", "一个人能描述清楚"],
                "考虑编排", ["能拆成独立子问题", "需要并行", "需要独立 / 对抗复核"], GREEN, YELLOW)
    rect(s, cx(0), 696, cw(12), 72, INK)
    text(s, cx(0) + 32, 712, cw(12) - 64, 40, "能单 Agent，就不编排。", 27, WHITE, "bold")
    out.append(finish(s))

    # 18
    s = page_base(18)
    label(s, cx(0), 256, "第三方社区插件 · 非 OpenCode 官方", RED)
    text(s, cx(0), 320, cw(4), 136, "长任务会带来\n上下文膨胀与注意力漂移", 28, INK, "bold")
    text(s, cx(0), 496, cw(4), 80, "复杂任务分职责；\n简单任务不要用。", 22, SLATE, "medium")
    roles = ["Orchestrator", "Explorer", "Oracle", "Council", "Librarian", "Designer", "Fixer", "Observer · 可选"]
    for i, name in enumerate(roles):
        x = cx(5 + (i % 2) * 4); y = 280 + (i // 2) * 88
        text(s, x, y, cw(3), 32, name, 18, BLUE if i != 3 else RED, "bold", font="latin")
        line(s, x, y + 48, x + cw(3), y + 48, RULE, 1)
    rect(s, cx(5), 664, cw(7), 72, "F8F4E8", YELLOW, 2)
    text(s, cx(5) + 24, 680, cw(7) - 48, 40, "“计划 → 实现 → 复核”是抽象阶段，不是角色名。", 19, INK, "bold")
    out.append(finish(s))

    # 19
    s = page_base(19)
    terminal(s, cx(0), 248, cw(8), 192, "shell · 课堂不执行", [
        ("$ npx oh-my-opencode-slim@latest install", TERMINAL_TEXT),
        ("→ 完成登录 / 刷新", CODE_MUTED),
        ("$ opencode", TERMINAL_TEXT),
    ])
    rect(s, cx(9), 248, cw(3), 192, WHITE, BLUE, 2)
    text(s, cx(9) + 24, 272, cw(3) - 48, 28, "OpenCode 对话", 17, BLUE, "bold")
    line(s, cx(9) + 24, 312, 1504 - 24, 312, RULE, 1)
    text(s, cx(9) + 24, 352, cw(3) - 48, 48, "> ping all agents", 19, INK, "bold", font="mono")
    label(s, cx(0), 472, "配置路径 · 安装不创建 agent 目录", RED)
    text(s, cx(0), 520, cw(12), 32, "用户级  ~/.config/opencode/oh-my-opencode-slim.json(c)", 18, INK, "bold", font="mono")
    text(s, cx(0), 560, cw(12), 32, "项目级  .opencode/oh-my-opencode-slim.json  · 可选覆盖", 18, INK, "bold", font="mono")
    label(s, cx(0), 624, "环境变量 · 后台编排 / 整体禁用")
    text(s, cx(0), 672, cw(7), 32, "OPENCODE_EXPERIMENTAL_BACKGROUND_SUBAGENTS=true", 15, BLUE, "bold", font="mono")
    text(s, cx(7), 672, cw(5), 32, "OH_MY_OPENCODE_SLIM_DISABLE=1", 15, RED, "bold", font="mono")
    text(s, cx(0), 728, cw(12), 32, "读写配置 / 调用模型需看清权限 · Council 多模型投票成本高", 18, SLATE, "medium")
    out.append(finish(s))

    # 20
    s = page_base(20)
    roles = [("01", "Orchestrator", "分派与汇总"), ("02", "Explorer", "只读侦察"),
             ("03", "Fixer", "有界实现"), ("04", "Oracle", "按需高风险评审"),
             ("05", "Council", "多模型投票 · 成本高")]
    for i, (n, name, job) in enumerate(roles):
        x = cx((i % 3) * 4); y = 264 + (i // 3) * 176
        text(s, x, y, cw(1), 32, n, 16, BLUE, "bold", font="latin")
        text(s, x, y + 48, cw(3), 36, name, 21, INK, "bold", font="latin")
        line(s, x, y + 96, x + cw(3), y + 96, RED if name == "Council" else BLUE, 2)
        text(s, x, y + 112, cw(3), 48, job, 19, SLATE, "medium")
    text(s, cx(8), 616, cw(4), 56, "Designer：UI\nLibrarian：官方资料", 18, SLATE, "medium")
    quote_mark(s, cx(0), 696, 56, GREEN)
    text(s, cx(0) + 32, 688, cw(9), 72, "按需启用；本课演示主要使用单 Agent。", 23, INK, "bold")
    out.append(finish(s))

    # 21
    s = page_base(21, dark=True)
    text(s, cx(0), 264, cw(7), 104, "把“符合需求”\n变成可核对的证据。", 39, WHITE, "bold")
    stages = ["需求", "规范", "实现", "核对"]
    for i, item in enumerate(stages):
        x = cx(i * 3)
        text(s, x, 488, cw(2), 48, item, 27, WHITE, "bold")
        if i < 3:
            line(s, x + cw(2), 512, cx((i + 1) * 3) - 24, 512, YELLOW, 2, True)
    text(s, cx(0), 648, cw(6), 32, "OpenSpec · openspec.dev", 18, "AFC0C2", font="latin")
    out.append(finish(s))

    # 22
    s = page_base(22)
    terminal(s, cx(0), 256, cw(6), 480, "openspec / directory", [
        ("openspec/", TERMINAL_TEXT),
        ("├─ changes/<change>/", TERMINAL_TEXT),
        ("│  ├─ proposal.md", CODE_MUTED),
        ("│  ├─ specs/.../spec.md", CODE_MUTED),
        ("│  ├─ design.md", CODE_MUTED),
        ("│  └─ tasks.md", CODE_MUTED),
        ("└─ changes/archive/...", TERMINAL_TEXT),
    ])
    artifacts = [("proposal", "为什么改"), ("spec", "SHALL + 场景"), ("design", "设计决策"), ("tasks", "最小步骤")]
    for i, (a, b) in enumerate(artifacts):
        y = 264 + i * 96
        text(s, cx(7), y, cw(2), 32, a, 20, BLUE, "bold", font="latin")
        text(s, cx(9), y, cw(3), 40, b, 22, INK, "medium")
    text(s, cx(7), 664, cw(5), 64, "培训夹具位于 docs 下，\n不进入 CLI 扫描。", 20, RED, "bold")
    out.append(finish(s))

    # 23
    s = page_base(23)
    two_columns(s, "CLI · 给人操作", ["终端中直接运行", "适合手动检查", "由人控制参数"],
                "聊天命令 · 给会话使用", ["本仓库使用连字符", "Agent 在会话内调用", "项目已有资产"], BLUE, BLUE)
    rect(s, cx(0), 696, cw(12), 72, TERMINAL)
    text(s, cx(0) + 24, 712, cw(12) - 48, 40,
         "/opsx-propose   /opsx-apply   /opsx-verify   /opsx-archive   /opsx-explore", 18,
         TERMINAL_TEXT, "bold", font="mono")
    out.append(finish(s))

    # 24
    s = page_base(24)
    stages = [("explore", "摸现状"), ("propose", "出规格"), ("apply", "按规格实现"), ("archive", "完成后归档")]
    for i, (a, b) in enumerate(stages):
        x = cx(i * 3)
        text(s, x, 280, cw(2), 36, a, 21, BLUE, "bold", font="latin")
        text(s, x, 336, cw(2), 48, b, 24, INK, "bold")
        if i < 3:
            line(s, x + cw(2), 360, cx((i + 1) * 3) - 24, 360, BLUE, 2, True)
    doc_excerpt(s, cx(0), 456, cw(7), 248, ".opencode/commands/opsx-verify.md", "项目已有启发式命令", [
        "读取 artifacts · 搜索代码",
        "判断完整性 / 正确性 / 连贯性",
        "无需 profile · 不自动运行测试",
    ])
    rect(s, cx(8), 456, cw(4), 248, "F7EDE8", RED, 2)
    text(s, cx(8) + 32, 488, cw(4) - 64, 32, "STRICT", 18, RED, "bold", font="latin")
    text(s, cx(8) + 32, 544, cw(4) - 64, 88, "格式 · requirement / scenario\n结构与可解析性", 20, INK, "bold")
    text(s, cx(8) + 32, 648, cw(4) - 64, 40, "不证明 artifacts 完成度或行为", 17, RED, "bold")
    out.append(finish(s))

    # 25
    s = page_base(25)
    text(s, cx(0), 280, cw(7), 112, "内置工具不够，\n又必须连接外部系统时，才考虑 MCP。", 34, INK, "bold")
    line(s, cx(0), 448, cx(7) - 32, 448, BLUE, 4)
    text(s, cx(0), 488, cw(7), 64, "它是协议，不是需要安装的单一运行时。", 24, SLATE, "medium")
    ruled_note(s, cx(8), 280, cw(4), 304, "本项目边界", "无 MCP 配置\n本课不演示", YELLOW)
    text(s, cx(8), 640, cw(4), 40, "别为用而用。", 24, RED, "bold")
    out.append(finish(s))

    # 26
    s = page_base(26, dark=True)
    line(s, cx(6) - 16, 256, cx(6) - 16, 760, "45525A", 1)
    label(s, cx(0), 264, "TDD · 已知目标行为", YELLOW)
    text(s, cx(0), 328, cw(5), 96, "“保存后应立即刷新”", 29, WHITE, "bold")
    text(s, cx(0), 464, cw(5), 112, "第一步\n写一个会失败的测试", 25, "C8D0CE", "medium")
    text(s, cx(0), 616, cw(5), 48, "项目案例完整演示", 20, "D9C982", "bold")
    label(s, cx(6), 264, "DIAGNOSIS · 未知故障原因", YELLOW)
    text(s, cx(6), 320, cw(6), 72, "虚构教学微案例\n搜索框 Enter 失败 · 点击正常", 24, WHITE, "bold")
    diagnosis = [
        ("01 复现 / 最小化", "只留搜索表单与 keydown 路径"),
        ("02 假设", "焦点 · preventDefault · submit 绑定"),
        ("03 插桩 / 排除", "keydown 已触发 · submit handler=0\n排除焦点 / 后端"),
        ("04 修复 / 回归", "统一 form onSubmit\nEnter / 点击各触发 1 次"),
    ]
    for i, (head, body) in enumerate(diagnosis):
        y = 416 + i * 80
        text(s, cx(6), y, cw(2), 32, head, 16, YELLOW, "bold")
        text(s, cx(8), y, cw(4), 56, body, 16, WHITE, "medium", leading=1.15)
    text(s, cx(6), 744, cw(6), 24, "独立于项目四阶段 · 不声称来自真实缺陷", 15, "D9C982", "bold")
    out.append(finish(s))

    # 27
    s = page_base(27)
    layers = [
        ("01", "strict", "规格结构完整", "不证明行为", BLUE, 88),
        ("02", "规范映射 / 启发式核对", "历史人工/等效规范映射\n/opsx-verify 可辅助", "当时未运行命令；不覆盖场景外", YELLOW, 104),
        ("03", "机器执行证据", "tests：被断言行为 · typecheck：类型一致\nbuild / package：可构建打包 · E2E：特定环境路径", "不证明未断言行为 / 其他环境 / 代码质量", GREEN, 144),
        ("04", "code review", "安全、清晰、可维护", "不穷尽所有行为", INK, 88),
    ]
    y = 248
    for n, head, can, cannot, color, row_h in layers:
        text(s, cx(0), y, cw(1), 32, n, 17, color, "bold", font="latin")
        text(s, cx(1), y, cw(3), 36, head, 21, INK, "bold", font="latin")
        text(s, cx(4), y, cw(4), row_h - 24, can, 18 if n == "03" else 20, GREEN if n == "03" else INK, "medium", leading=1.2)
        text(s, cx(8), y, cw(4), row_h - 24, cannot, 18, RED, "medium", leading=1.2)
        line(s, cx(0), y + row_h - 8, 1504, y + row_h - 8, RULE, 1)
        y += row_h
    text(s, cx(0), 760, cw(12), 32, "四层互补；机器结果不能读成“全部正确”。", 21, BLUE, "bold")
    out.append(finish(s))

    # 28
    s = page_base(28, dark=True)
    text(s, cx(0), 264, cw(7), 64, "原始需求", 20, YELLOW, "bold")
    text(s, cx(0), 336, cw(8), 160, "“在提醒面板里让负责人\n配置临期窗口。”", 37, WHITE, "bold")
    ruled_note(s, cx(8), 272, cw(4), 320, "业务背景", "临期\n今日到期\n已逾期", YELLOW)
    text(s, cx(0), 632, cw(8), 80, "一句话里，藏着五个可验证的歧义。", 25, "C7D0CE", "medium")
    text(s, cx(8), 640, cw(4), 64, "来自项目 · 范围小\n可观察 · 已脱敏", 18, "A8B6B5")
    out.append(finish(s))

    # 29
    s = page_base(29)
    questions = [
        ("Q1", "0 合法吗？", "合法；未来提醒不标临期"),
        ("Q2", "能填多大？", "0..9007199254740991"),
        ("Q3", "何时生效？", "显式保存后刷新并持久化"),
    ]
    for i, (q, ask, answer) in enumerate(questions):
        y = 256 + i * 136
        text(s, cx(0), y, cw(1), 32, q, 18, YELLOW, "bold", font="latin")
        text(s, cx(1), y, cw(4), 40, ask, 25, INK, "bold")
        line(s, cx(5), y + 18, cx(7) - 32, y + 18, RULE, 1)
        text(s, cx(7), y, cw(5), 64, answer, 22, BLUE, "bold")
    quote_mark(s, cx(0), 680, 56, GREEN)
    text(s, cx(0) + 32, 672, cw(11), 72, "保存前不生效 · 保存后立即刷新 · 重开保持 · 提醒日期可见", 22, INK, "bold")
    out.append(finish(s))

    # 30
    s = page_base(30)
    doc_excerpt(s, cx(0), 256, cw(7), 472,
                "docs/training/vibe-coding/training-change/spec.md",
                "Scenario: 临期窗口为 0", [
                    "GIVEN 配置为 0 并显式保存",
                    "THEN 未来提醒不归入临期",
                    "AND 今日到期 / 已逾期不受影响",
                ], GREEN)
    label(s, cx(8), 264, "归档漂移 · 三条事实", RED)
    facts = [
        "原 change 已归档",
        "正式 spec 仍是旧版",
        "最终实现未合入 main",
    ]
    for i, fact in enumerate(facts):
        text(s, cx(8), 328 + i * 88, cw(4), 48, fact, 22, INK, "bold")
        line(s, cx(8), 392 + i * 88, 1504, 392 + i * 88, RULE, 1)
    text(s, cx(8), 600, cw(4), 80, "openspec/changes/archive/\n2026-08-10-add-relocation-\nservice-workbench/", 14, SLATE, font="mono", leading=1.15)
    text(s, cx(8), 696, cw(4), 40, "培训夹具，无产品批准语义", 19, RED, "bold")
    out.append(finish(s))

    # 31
    s = page_base(31)
    text(s, cx(0), 256, cw(6), 32, "RED-TEST · f771667", 18, RED, "bold", font="latin")
    terminal(s, cx(0), 304, cw(6), 336, "historical output", [
        ("domain + integration   32 passed", "B9D7C0"),
        ("main                  2 failed", "E8A6A0"),
        ("V2_MUTATION_UNKNOWN", "E8A6A0"),
        ("renderer              1 failed", "E8A6A0"),
        ("未找到“临期窗口”控件", "E8A6A0"),
    ], group="evidence-pair")
    text(s, cx(6), 256, cw(6), 32, "PRESET-FAILURE · 0554164", 18, RED, "bold", font="latin")
    terminal(s, cx(6), 304, cw(6), 336, "historical output", [
        ("IPC / facade / UI     connected", "B9D7C0"),
        ("0 case                failed", "E8A6A0"),
        ("“必须不小于 1”", "E8A6A0"),
        ("negative / fraction   rejected", "B9D7C0"),
        ("renderer: mutation not called", "E8A6A0"),
    ], group="evidence-pair")
    text(s, cx(0), 696, cw(12), 56, "红：行为已经写明，实现还没跟上。表面能点，也可能做错。", 24, INK, "bold")
    out.append(finish(s))

    # 32
    s = page_base(32)
    terminal(s, cx(0), 256, cw(5), 320, "green-final · d304bd5 · historical", [
        ("focused       95 passed", "B9D7C0"),
        ("full          1076 passed", "B9D7C0"),
        ("strict        valid (structure)", "B9D7C0"),
        ("e2e build     passed", "B9D7C0"),
        ("layout E2E    1 / 1", "B9D7C0"),
    ])
    doc_excerpt(s, cx(6), 256, cw(6), 320, "review / historical ledger", "质量门", [
        "2 个 medium 已修正；2 个既有 low 保留",
        "最后三项修正后未做第三轮独立 Oracle 评审",
        "当前是课程内容终审",
    ], GREEN)
    rect(s, cx(0), 608, cw(7), 120, "E8EEE9", GREEN, 2)
    text(s, cx(0) + 24, 624, cw(7) - 48, 88, "UI：填 0 → 保存 → 立即刷新 → 日期可见\n今日到期 / 已逾期不变", 20, INK, "bold", leading=1.15)
    rect(s, cx(8), 608, cw(4), 120, "F7EDE8", RED, 2)
    text(s, cx(8) + 24, 624, cw(4) - 48, 88, "strict：2026-08-10 历史摘要\n规范映射：历史人工 / 等效核对\n当时未运行 /opsx-verify", 14, RED, "bold", leading=1.12)
    out.append(finish(s))

    # 33
    s = page_base(33)
    checks = ["先文档后提问", "写清约定与证据", "TDD 先红", "diff 把关再绿", "review 看质量", "verify 对约定"]
    for i, item in enumerate(checks):
        x = cx((i % 3) * 4); y = 256 + (i // 3) * 120
        text(s, x, y, cw(1), 32, f"{i+1:02d}", 17, BLUE, "bold", font="latin")
        text(s, x, y + 40, cw(3), 48, item, 22, INK, "bold")
    line(s, cx(0), 520, 1504, 520, RULE, 1)
    risks = [("低", "简短约定 + 聚焦验证"), ("中", "加深文档 / 测试 / review"), ("高", "人工确认 + 完整证据")]
    for i, (level, action) in enumerate(risks):
        y = 560 + i * 64
        text(s, cx(0), y, cw(1), 32, level, 20, RED if level == "高" else BLUE, "bold")
        text(s, cx(2), y, cw(7), 32, action, 21, INK, "medium")
    text(s, cx(9), 560, cw(3), 128, "目标、边界、证据\n任何强度都不能省。", 23, GREEN, "bold", align="right")
    text(s, cx(0), 752, cw(8), 24, "跨 shared IPC / main / renderer：集成风险按中等", 16, SLATE, font="mono")
    out.append(finish(s))

    # 34
    s = page_base(34, dark=True)
    text(s, 96, 88, cw(3), 32, "Q&A", 21, YELLOW, "bold", font="latin")
    text(s, 96, 184, cw(8), 80, "先答共性，再记个性。", 43, WHITE, "bold")
    text(s, 96, 312, cw(7), 64, "个人配置与扩展问题进入待答问题池，\n会后书面处理。", 25, "C1CBCA", "medium")
    line(s, 96, 456, 1504, 456, "46535A", 1)
    words = ["受控", "证据", "闭环"]
    for i, word in enumerate(words):
        text(s, cx(i * 4), 528, cw(3), 72, word, 35, WHITE, "bold")
        line(s, cx(i * 4), 616, cx(i * 4) + cw(3), 616, YELLOW if i == 1 else "59666C", 3)
    text(s, 96, 712, cw(10), 64, "明天开始：先写目标、边界和验收证据，再打开会话。", 24, "D6DEDB", "bold")
    text(s, 1456, 840, 48, 20, "34", 13, "91A0A5", align="right", font="latin")
    out.append(s)

    if len(out) != 34:
        raise ValueError(f"Expected 34 slides, got {len(out)}")
    return out


def rgb(value: str) -> tuple[int, int, int]:
    return cast(tuple[int, int, int], tuple(int(value[i:i + 2], 16) for i in (0, 2, 4)))


def font_candidates() -> list[str]:
    env = os.environ.get("VIBE_SLIDES_CJK_FONT")
    return [p for p in [
        env,
        "/System/Library/AssetsV2/com_apple_MobileAsset_Font8/86ba2c91f017a3749571a82f2c6d890ac7ffb2fb.asset/AssetData/PingFang.ttc",
        "/System/Library/Fonts/PingFang.ttc",
        "/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc",
        "/usr/share/fonts/opentype/noto/NotoSansCJKsc-Regular.otf",
        str(Path.home() / "Library/Fonts/SourceHanSansSC-Regular.otf"),
    ] if p]


def resolve_font() -> Path:
    for item in font_candidates():
        path = Path(item).expanduser()
        if path.is_file():
            return path
    raise FileNotFoundError("Set VIBE_SLIDES_CJK_FONT to a Chinese .ttf/.otf/.ttc font")


def font_indexes(path: Path) -> dict[str, int]:
    if path.suffix.lower() != ".ttc":
        return {"regular": 0, "medium": 0, "bold": 0}
    return {
        "regular": int(os.environ.get("VIBE_SLIDES_CJK_REGULAR_INDEX", "3")),
        "medium": int(os.environ.get("VIBE_SLIDES_CJK_MEDIUM_INDEX", "7")),
        "bold": int(os.environ.get("VIBE_SLIDES_CJK_BOLD_INDEX", "11")),
    }


def pil_font(path: Path, size: int, weight: str) -> ImageFont.FreeTypeFont:
    return ImageFont.truetype(str(path), max(1, size), index=font_indexes(path)[weight])


def render_pptx(slides: Iterable[Slide], output: Path) -> None:
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(PPT_W), Inches(PPT_H)
    sx, sy = PPT_W / W, PPT_H / H
    for model in slides:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        for e in model.elements:
            p = e.props
            x, y, w, h = Inches(e.x * sx), Inches(e.y * sy), Inches(e.w * sx), Inches(e.h * sy)
            if e.kind == "rect":
                shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if p.get("radius") else MSO_SHAPE.RECTANGLE
                shape = slide.shapes.add_shape(shape_type, x, y, w, h)
                shape.fill.solid(); shape.fill.fore_color.rgb = RGBColor(*rgb(p["fill"]))
                if p.get("stroke"):
                    shape.line.color.rgb = RGBColor(*rgb(p["stroke"])); shape.line.width = Pt(p.get("sw", 1) * .75)
                else:
                    shape.line.fill.background()
            elif e.kind == "line":
                end_x, end_y = Inches((e.x + e.w) * sx), Inches((e.y + e.h) * sy)
                shape = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x, y, end_x, end_y)
                shape.line.color.rgb = RGBColor(*rgb(p["color"])); shape.line.width = Pt(p.get("sw", 2) * .75)
                if p.get("arrow"):
                    marker = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE,
                        Inches((e.x + e.w) * sx - .045), Inches((e.y + e.h) * sy - .045), Inches(.09), Inches(.09))
                    marker.fill.solid(); marker.fill.fore_color.rgb = RGBColor(*rgb(p["color"])); marker.line.fill.background()
                    marker.rotation = math.degrees(math.atan2(e.h, e.w)) + 90
            elif e.kind == "text":
                box = slide.shapes.add_textbox(x, y, w, h)
                tf = box.text_frame; tf.clear(); tf.word_wrap = True
                tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
                tf.vertical_anchor = {"top": MSO_ANCHOR.TOP, "middle": MSO_ANCHOR.MIDDLE, "bottom": MSO_ANCHOR.BOTTOM}[p["valign"]]
                for i, value in enumerate(p["text"].split("\n")):
                    para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                    para.text = value; para.space_after = Pt(0); para.line_spacing = p["leading"]
                    para.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}[p["align"]]
                    run = para.runs[0] if para.runs else para.add_run()
                    run.font.size = Pt(p["size"] * .75); run.font.bold = p["weight"] == "bold"
                    run.font.color.rgb = RGBColor(*rgb(p["color"]))
                    family = "Aptos Mono" if p["font"] == "mono" else ("Aptos" if p["font"] == "latin" else "Microsoft YaHei")
                    run.font.name = family
                    rpr = run._r.get_or_add_rPr(); rpr.set(qn("a:ea"), "Microsoft YaHei"); rpr.set(qn("a:latin"), family)
    prs.save(output)


def draw_arrow(draw: ImageDraw.ImageDraw, a: tuple[int, int], b: tuple[int, int], color: tuple[int, int, int], size: int = 12) -> None:
    angle = math.atan2(b[1] - a[1], b[0] - a[0])
    draw.polygon([b, (b[0] - size * math.cos(angle - .5), b[1] - size * math.sin(angle - .5)),
                  (b[0] - size * math.cos(angle + .5), b[1] - size * math.sin(angle + .5))], fill=color)


def render_image(slide: Slide, font_path: Path, scale: float = 1.0) -> Image.Image:
    image = Image.new("RGB", (int(W * scale), int(H * scale)), rgb(INK if slide.dark else PAPER))
    draw = ImageDraw.Draw(image)
    for e in slide.elements:
        p = e.props
        box = tuple(int(v * scale) for v in (e.x, e.y, e.x + e.w, e.y + e.h))
        if e.kind == "rect":
            draw.rounded_rectangle(box, radius=int(p.get("radius", 0) * scale), fill=rgb(p["fill"]),
                                   outline=rgb(p["stroke"]) if p.get("stroke") else None,
                                   width=max(1, int(p.get("sw", 1) * scale)))
        elif e.kind == "line":
            a = (int(e.x * scale), int(e.y * scale)); b = (int((e.x + e.w) * scale), int((e.y + e.h) * scale))
            draw.line([a, b], fill=rgb(p["color"]), width=max(1, int(p.get("sw", 2) * scale)))
            if p.get("arrow"): draw_arrow(draw, a, b, rgb(p["color"]), max(8, int(12 * scale)))
        elif e.kind == "text":
            font = pil_font(font_path, int(p["size"] * scale), p["weight"])
            spacing = int(p["size"] * (p["leading"] - 1) * scale)
            bbox = draw.multiline_textbbox((0, 0), p["text"], font=font, spacing=spacing, align=p["align"])
            tw, th = bbox[2] - bbox[0], bbox[3] - bbox[1]
            tx = e.x * scale if p["align"] == "left" else ((e.x + e.w / 2) * scale - tw / 2 if p["align"] == "center" else (e.x + e.w) * scale - tw)
            ty = e.y * scale if p["valign"] == "top" else ((e.y + e.h / 2) * scale - th / 2 if p["valign"] == "middle" else (e.y + e.h) * scale - th)
            draw.multiline_text((tx, ty - bbox[1]), p["text"], font=font, fill=rgb(p["color"]), spacing=spacing, align=p["align"])
    return image


def render_pdf(slides: Iterable[Slide], output: Path, font_path: Path) -> None:
    c = canvas.Canvas(str(output), pagesize=(960, 540), pageCompression=1)
    for slide in slides:
        image = render_image(slide, font_path)
        buffer = BytesIO(); image.save(buffer, "PNG", optimize=True); buffer.seek(0)
        c.drawImage(ImageReader(buffer), 0, 0, 960, 540, preserveAspectRatio=True)
        c.showPage()
    c.save()


def contact_sheet(paths: list[Path], output: Path) -> None:
    tw, th, cols = 320, 180, 4
    rows = math.ceil(len(paths) / cols)
    sheet = Image.new("RGB", (tw * cols, th * rows), rgb(INK))
    for i, path in enumerate(paths):
        im = Image.open(path).convert("RGB"); im.thumbnail((tw, th), Image.Resampling.LANCZOS)
        sheet.paste(im, ((i % cols) * tw, (i // cols) * th))
    sheet.save(output, optimize=True)


def render_previews(slides: Iterable[Slide], output: Path, font_path: Path) -> None:
    output.mkdir(parents=True, exist_ok=True)
    paths = []
    for slide in slides:
        path = output / f"slide-{slide.number:02d}.png"
        render_image(slide, font_path).save(path, optimize=True)
        paths.append(path)
    contact_sheet(paths, output / "contact-sheet.png")


def validate(slides: list[Slide], font_path: Path) -> list[str]:
    issues: list[str] = []
    probe = ImageDraw.Draw(Image.new("RGB", (8, 8)))
    if [s.title for s in slides] != TITLES:
        issues.append("title-order mismatch")
    for slide in slides:
        for e in slide.elements:
            if e.x < 0 or e.y < 0 or e.x + e.w > W or e.y + e.h > H:
                issues.append(f"p{slide.number}: out of bounds {e.kind}")
            if e.kind == "text":
                p = e.props; font = pil_font(font_path, int(p["size"]), p["weight"])
                spacing = int(p["size"] * (p["leading"] - 1))
                box = probe.multiline_textbbox((0, 0), p["text"], font=font, spacing=spacing, align=p["align"])
                if box[2] - box[0] > e.w + 3 or box[3] - box[1] > e.h + 3:
                    issues.append(f"p{slide.number}: text overflow {p['text']!r}")
                if round(e.y) % 8:
                    issues.append(f"p{slide.number}: text baseline off-grid at y={e.y}")
        groups: dict[str, list[Element]] = {}
        for e in slide.elements:
            if e.props.get("group"):
                groups.setdefault(e.props["group"], []).append(e)
        for name, members in groups.items():
            if len({round(e.y, 1) for e in members}) > 1:
                issues.append(f"p{slide.number}: group {name} top alignment")
    return issues


def main() -> None:
    root = Path(__file__).resolve().parent
    slides = make_slides(); font_path = resolve_font()
    issues = validate(slides, font_path)
    if issues:
        raise ValueError("Layout validation failed:\n" + "\n".join(issues))
    render_pptx(slides, root / "vibe-coding-training.pptx")
    render_pdf(slides, root / "vibe-coding-training.pdf", font_path)
    if preview := os.environ.get("VIBE_SLIDES_PREVIEW_DIR"):
        render_previews(slides, Path(preview).expanduser(), font_path)
    print(f"Generated {len(slides)} pages")
    print(f"CJK render font: {font_path}")


if __name__ == "__main__":
    main()
